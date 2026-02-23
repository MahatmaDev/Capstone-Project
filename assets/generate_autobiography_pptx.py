from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
import os, re

# =========================
# CONFIG — EDIT THESE
# =========================
DOCX_PATH = r"D:\Capstone\assets\documents\Kelvin_Ndegwa_Autobiography.docx"
LOGO_PATH = r"D:\Capstone\assets\images\DeKUT.png"   # put logo here
OUT_PATH  = r"D:\Capstone\assets\presentations\Autobiography_Presentation_Final_DeKUT.pptx"

NAME = "Kelvin Maina Ndegwa"
REG_NO = "C025-01-2566/2022"
COURSE = "Bachelor of Science in Information Technology"
INSTITUTION = "Dedan Kimathi University of Technology"
TOPIC = "Autobiography Presentation"

FONT = "Times New Roman"
TITLE_SIZE = Pt(34)
SUBTITLE_SIZE = Pt(18)
BODY_SIZE = Pt(20)
LINE_SPACING = 1.5

# =========================
# PPT HELPERS
# =========================
def set_run(run, size=BODY_SIZE, bold=False, color=(0,0,0)):
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)

def set_white_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

def add_logo(slide, prs):
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(
            LOGO_PATH,
            prs.slide_width - Inches(1.25),
            Inches(0.2),
            width=Inches(1.0)
        )

def add_slide_number(slide, prs, num):
    box = slide.shapes.add_textbox(
        prs.slide_width - Inches(1.0),
        prs.slide_height - Inches(0.5),
        Inches(0.8),
        Inches(0.3)
    )
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(num)
    p.alignment = PP_ALIGN.RIGHT
    r = p.runs[0]
    set_run(r, size=Pt(12), bold=False, color=(120,120,120))

def format_paragraph(p, align=PP_ALIGN.JUSTIFY):
    p.alignment = align
    p.line_spacing = LINE_SPACING
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    for r in p.runs:
        r.font.name = FONT

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_white_bg(slide)

    tf = slide.shapes.title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = TOPIC
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    set_run(p.runs[0], size=TITLE_SIZE, bold=True)

    sub = slide.placeholders[1].text_frame
    sub.clear()
    lines = [
        f"Name: {NAME}",
        f"Reg No: {REG_NO}",
        f"Course: {COURSE}",
        f"Institution: {INSTITUTION}",
        "Date: ____________________",
    ]
    for i, line in enumerate(lines):
        p = sub.paragraphs[0] if i == 0 else sub.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.0
        set_run(p.runs[0], size=SUBTITLE_SIZE, bold=False)

    add_logo(slide, prs)

def add_bullets_slide(prs, title, bullets, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_white_bg(slide)

    slide.shapes.title.text = title
    ttf = slide.shapes.title.text_frame
    for p in ttf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.0
        for r in p.runs:
            set_run(r, size=Pt(28), bold=True)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        format_paragraph(p, align=PP_ALIGN.JUSTIFY)
        set_run(p.runs[0], size=BODY_SIZE, bold=False)

    add_logo(slide, prs)
    add_slide_number(slide, prs, slide_num)

# =========================
# DOCX → SECTIONS
# =========================
def is_heading(line: str) -> bool:
    if len(line) > 90:
        return False
    if line.isupper() and len(line.split()) <= 10:
        return True
    # Common section starters
    starters = (
        "introduction", "early life", "early education", "secondary education",
        "university", "professional certifications", "technical skills",
        "creative journey", "digital presence", "current project",
        "vibewave", "challenges", "values", "present life",
        "future aspirations", "industrial attachment", "conclusion"
    )
    low = line.lower()
    return any(low.startswith(s) for s in starters)

def summarize_to_bullets(paragraphs, max_bullets=6):
    text = " ".join(paragraphs)
    sentences = re.split(r"(?<=[.!?])\s+", text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 25]

    bullets = []
    for s in sentences:
        if len(s) > 150:
            s = s[:147].rsplit(" ", 1)[0] + "..."
        bullets.append(s)
        if len(bullets) >= max_bullets:
            break

    if not bullets:
        bullets = [p[:147] + "..." if len(p) > 150 else p for p in paragraphs[:max_bullets]]

    return bullets[:max_bullets]

def extract_sections(docx_path):
    doc = Document(docx_path)
    paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

    sections = []
    current = {"title": "Overview", "content": []}

    for t in paras:
        if is_heading(t):
            if current["content"]:
                sections.append(current)
            current = {"title": t, "content": []}
        else:
            current["content"].append(t)

    if current["content"]:
        sections.append(current)

    slides = []
    for sec in sections:
        bullets = summarize_to_bullets(sec["content"], max_bullets=6)
        slides.append((sec["title"], bullets))

    return slides

def main():
    if not os.path.exists(DOCX_PATH):
        print("ERROR: DOCX not found:", DOCX_PATH)
        return

    prs = Presentation()
    add_title_slide(prs)

    slides = extract_sections(DOCX_PATH)

    # Outline slide
    outline = [t for t, _ in slides][:10]
    add_bullets_slide(prs, "Outline", outline, slide_num=1)

    # Add content slides (panel friendly: 12–16 slides)
    slide_no = 2
    for title, bullets in slides[:14]:
        add_bullets_slide(prs, title, bullets, slide_num=slide_no)
        slide_no += 1

    prs.save(OUT_PATH)
    print("DONE:", OUT_PATH)

if __name__ == "__main__":
    main()