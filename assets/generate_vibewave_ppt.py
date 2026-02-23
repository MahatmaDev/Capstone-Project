from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# =============================
# SCHOOL DETAILS (Edit anytime)
# =============================
NAME = "Kelvin Maina Ndegwa"
REG_NO = "C025-01-2566/2022"
COURSE = "Bachelor of Science in Information Technology"
INSTITUTION = "Dedan Kimathi University of Technology"
TOPIC = "VibeWave Innovation"
DATE_LINE = "Date: ____________________"

# =============================
# FORMATTING (School format)
# =============================
FONT = "Times New Roman"
BODY_SIZE = Pt(22)
TITLE_SIZE = Pt(34)
SUBTITLE_SIZE = Pt(20)
LINE_SPACING = 1.5

def set_run(run, size=BODY_SIZE, bold=False):
    run.font.name = FONT
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = RGBColor(0, 0, 0)

def format_paragraph(p, align=PP_ALIGN.JUSTIFY):
    p.alignment = align
    p.line_spacing = LINE_SPACING
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    for r in p.runs:
        r.font.name = FONT

def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    tf = slide.shapes.title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = TOPIC
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    set_run(p.runs[0], size=TITLE_SIZE, bold=True)

    sub = slide.placeholders[1].text_frame
    sub.clear()
    lines = [
        f"Name: {NAME}",
        f"Reg No: {REG_NO}",
        f"Course: {COURSE}",
        f"Institution: {INSTITUTION}",
        DATE_LINE,
    ]
    for i, line in enumerate(lines):
        p = sub.paragraphs[0] if i == 0 else sub.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.0
        set_run(p.runs[0], size=SUBTITLE_SIZE, bold=False)

def add_bullets_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    # title formatting
    ttf = slide.shapes.title.text_frame
    for p in ttf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.0
        for r in p.runs:
            set_run(r, size=Pt(28), bold=True)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
        p.level = level
        format_paragraph(p, align=PP_ALIGN.JUSTIFY)
        set_run(p.runs[0], size=BODY_SIZE, bold=False)

def build():
    prs = Presentation()
    add_title_slide(prs)

    add_bullets_slide(prs, "Presentation Overview", [
        "Background & problem statement",
        "Goals and objectives",
        "Proposed solution (VibeWave)",
        "Key features and workflow",
        "Technology stack & architecture",
        "Implementation plan",
        "Benefits, impact & sustainability",
        "Risks, ethics & data privacy",
        "Conclusion",
    ])

    add_bullets_slide(prs, "Background & Problem Statement", [
        "Live entertainment increasingly demands interactive engagement, yet DJ–audience communication remains informal and chaotic.",
        "Music requests are unmanaged (shouted, handwritten, relayed), leading to delays and conflict.",
        "Payments/tips are fragmented with limited transparency and weak linkage to requests.",
        "Organizers lack engagement data (popular requests, peak moments, audience sentiment).",
        "These gaps reduce user experience, efficiency, and revenue opportunities.",
    ])

    add_bullets_slide(prs, "Goal & Objectives", [
        "Goal: Build a scalable digital platform that structures DJ–fan engagement during live events.",
        ("Enable real-time song requests through a mobile/web interface.", 1),
        ("Support secure digital payments linked to requests and tips.", 1),
        ("Provide voting/priority mechanisms to fairly manage requests.", 1),
        ("Offer analytics and insights for DJs, organizers and venues.", 1),
        ("Improve user experience through transparency and feedback.", 1),
    ])

    add_bullets_slide(prs, "Proposed Solution: VibeWave", [
        "VibeWave is a DJ–fan engagement platform integrating event interaction, payments, and engagement insights.",
        "It converts unstructured requests into a managed queue with configurable voting/priority rules.",
        "Fans interact via a simple interface; DJs manage requests via a DJ dashboard.",
        "Payments can be attached to requests (tips/priority) while maintaining fairness controls.",
    ])

    add_bullets_slide(prs, "Key Features (Summary)", [
        "Fan: join event, request songs, vote, track status, tip, rate the DJ.",
        "DJ: manage queue, moderate content, track payments, view analytics, configure event rules.",
        "Organizer: monitor engagement, improve planning using insights, ensure structured flow.",
    ])

    add_bullets_slide(prs, "System Workflow (High Level)", [
        "1) Organizer creates an event and sets rules (limits, voting, moderation).",
        "2) Fans join via QR/event code and submit requests.",
        "3) Requests are validated, deduplicated, and added to a queue.",
        "4) Fans vote; the system updates ranking/priority based on rules.",
        "5) DJ reviews the queue and plays songs; statuses update to users.",
        "6) Post-event feedback is collected; analytics are generated.",
    ])

    add_bullets_slide(prs, "Technology Stack", [
        ("Frontend / Client", 0),
        ("Android: Kotlin", 1),
        ("Web: HTML • CSS • JavaScript", 1),
        ("UX: Responsive, accessible design", 1),
        ("Backend / Cloud", 0),
        ("Firebase Authentication", 1),
        ("Firestore / Realtime Database", 1),
        ("Cloud Functions (serverless logic)", 1),
        ("Payments: Mobile money gateway integration (planned)", 1),
    ])

    add_bullets_slide(prs, "Architecture (Conceptual)", [
        "Client Layer: Fan app + DJ dashboard (web/mobile).",
        "Service Layer: authentication, validation, voting/ranking, notifications.",
        "Data Layer: events, requests, votes, payments, feedback.",
        "Analytics Layer: aggregations, dashboards, exports.",
        "Security: role-based access, moderation policies, audit logging.",
    ])

    add_bullets_slide(prs, "Implementation Plan", [
        "Phase 1: Requirements gathering and UI/UX design.",
        "Phase 2: MVP build (auth, event join, request submission, DJ queue).",
        "Phase 3: Voting, moderation, and real-time updates.",
        "Phase 4: Payments integration and transaction records.",
        "Phase 5: Analytics, testing, deployment, documentation.",
    ])

    add_bullets_slide(prs, "Benefits & Impact", [
        "Enhances live event experience through transparency and structured engagement.",
        "Creates revenue opportunities (tips, premium features, partnerships).",
        "Enables data-driven planning through engagement analytics.",
        "Accelerates digital transformation in entertainment.",
        "Promotes youth entrepreneurship and innovation.",
    ])

    add_bullets_slide(prs, "Risks, Ethics & Data Privacy", [
        "Data privacy: minimize personal data; enforce strong access controls.",
        "Fairness: prevent pay-to-win abuse through rules and moderation.",
        "Safety: filter harmful content and reduce spam/harassment.",
        "Compliance: align with relevant data protection requirements.",
        "Reliability: design for poor network conditions and safe fallback behavior.",
    ])

    add_bullets_slide(prs, "Conclusion", [
        "VibeWave modernizes DJ–fan interaction by combining structured requests, engagement tools and payments.",
        "It demonstrates practical application of IT to solve real-world challenges in entertainment.",
        "Next steps: finalize MVP, pilot at live events, iterate using feedback, and scale.",
        "Thank you.",
    ])

    out_file = "VibeWave_Innovation_Presentation_Capstone_Project.pptx"
    prs.save(out_file)
    print(f"Saved: {out_file}")

if __name__ == "__main__":
    build()